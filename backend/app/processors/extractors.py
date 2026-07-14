"""Optional-dependency extractors for local KnowledgeOS source formats."""

from __future__ import annotations

from pathlib import Path


class ExtractionError(RuntimeError):
    """Raised when a local processor is unavailable or cannot read a source."""


def extract_text(path: Path) -> list[tuple[str, str | None]]:
    """Return `(text, location)` units from a supported local source file."""

    extension = path.suffix.lower()
    if extension in {".txt", ".md", ".markdown"}:
        return [(path.read_text(encoding="utf-8", errors="replace"), None)]
    if extension == ".pdf":
        return _extract_pdf(path)
    if extension == ".docx":
        return _extract_docx(path)
    if extension == ".pptx":
        return _extract_pptx(path)
    if extension in {".png", ".jpg", ".jpeg", ".webp"}:
        return _extract_image(path)
    if extension in {".mp3", ".m4a", ".wav", ".mp4", ".mkv", ".avi", ".mov"}:
        return _extract_media(path)
    raise ExtractionError(f"No extractor is configured for {extension}")


def _extract_pdf(path: Path) -> list[tuple[str, str | None]]:
    try:
        from pypdf import PdfReader
    except ImportError as error:
        raise ExtractionError("Install pypdf to index PDF files") from error
    return [
        (page.extract_text() or "", f"page {number}")
        for number, page in enumerate(PdfReader(path).pages, start=1)
    ]


def _extract_docx(path: Path) -> list[tuple[str, str | None]]:
    try:
        from docx import Document
    except ImportError as error:
        raise ExtractionError("Install python-docx to index DOCX files") from error
    text = "\n".join(paragraph.text for paragraph in Document(path).paragraphs)
    return [(text, None)]


def _extract_pptx(path: Path) -> list[tuple[str, str | None]]:
    try:
        from pptx import Presentation
    except ImportError as error:
        raise ExtractionError("Install python-pptx to index PPTX files") from error
    return [
        (
            "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text")),
            f"slide {number}",
        )
        for number, slide in enumerate(Presentation(path).slides, start=1)
    ]


def _extract_image(path: Path) -> list[tuple[str, str | None]]:
    try:
        import pytesseract
        from PIL import Image
    except ImportError as error:
        raise ExtractionError("Install pytesseract and Pillow to OCR images") from error
    return [(pytesseract.image_to_string(Image.open(path)), None)]


def _extract_media(path: Path) -> list[tuple[str, str | None]]:
    try:
        import whisper
    except ImportError as error:
        raise ExtractionError("Install openai-whisper and FFmpeg to transcribe media") from error
    result = whisper.load_model("base").transcribe(str(path))
    return [
        (segment["text"].strip(), f"{segment['start']:.1f}s–{segment['end']:.1f}s")
        for segment in result["segments"]
    ]
